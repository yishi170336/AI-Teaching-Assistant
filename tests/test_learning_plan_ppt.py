from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from backend.app.services.learning_plan_ppt import (
    build_learning_plan_ppt,
    generate_learning_plan_ppt,
    parse_learning_plan,
    presentation_filename,
)


SAMPLE_PLAN = r"""
# 晶体管放大电路学习规划

**总体目标：** 建立从静态工作点到动态参数分析的完整知识链，能够独立分析基本共射放大电路。

### 学习现状与执行原则

- 当前基础：已经掌握欧姆定律，但对晶体管工作区与等效模型不熟悉。
- 先补前置，再进入参数计算；每个阶段都通过练习与自测收口。

## 阶段一：诊断与前置补全（1～2 小时）

- **目标：** 确认直流通路、交流通路和晶体管三个工作区的理解缺口。
- **具体行动：** 画出典型共射电路的直流通路与交流通路。
- 回顾欧姆定律、基尔霍夫定律以及电容的隔直通交作用。
- 完成 5 道前置自测题，并记录不会的知识点。
- **完成标准：** 能不看资料说明截止区、放大区和饱和区的判据。
- **资料依据：** [资料1] 基本放大电路工作原理。

## 阶段二：静态工作点分析（2～3 小时）

- **目标：** 掌握 $I_{BQ}$、$I_{CQ}$ 与 $U_{CEQ}$ 的计算路径。
- **具体行动：** 整理固定偏置和分压偏置两类电路的计算模板。
- 独立完成 6 道静态分析题，对每个结果标注单位。
- 用负载线检查计算结果是否处于放大区。
- **完成标准：** 静态参数计算正确率达到 85%，并能解释 Q 点偏移的影响。

## 阶段三：动态参数与等效模型（3～4 小时）

- **目标：** 使用 h 参数等效模型求解电压增益、输入电阻和输出电阻。
- **具体行动：** 先画交流等效电路，再列出输入与输出回路方程。
- 对比旁路电容存在与不存在时增益和输入电阻的变化。
- 推导 $A_u=-\frac{\beta R'_L}{r_{be}}$，并解释各参数对电压增益的影响。
- 完成 4 道综合计算题，并用量纲与数量级复核答案。
- **完成标准：** 能独立画出等效电路，三项动态参数正确率不低于 80%。

## 阶段四：专项训练与复盘（2 小时）

- **目标：** 把静态与动态分析串联起来，形成稳定的解题流程。
- **具体行动：** 完成一组限时综合题，按“识图—静态—动态—校验”书写。
- 将错题按概念、建模、计算和单位四类归因。
- 24 小时后重做错题，并用自己的话讲解完整思路。
- **完成标准：** 综合题正确率达到 80%，同类错误不连续出现两次。

### 7 天学习安排

- Day 1：完成诊断与前置知识补全。
- Day 2～3：集中训练静态工作点计算。
- Day 4～5：学习等效模型并完成动态参数练习。
- Day 6：完成综合题与错题归因。
- Day 7：闭卷自测、复盘并形成一页知识地图。

### 量化验收指标

- 能在 3 分钟内画出直流通路与交流通路。
- 静态工作点计算正确率达到 85%。
- 动态参数综合题正确率达到 80%。
- 能口头解释失真类型、产生原因与调整方向。
""".strip()


def test_parse_learning_plan_builds_a_logical_story() -> None:
    plan = parse_learning_plan(SAMPLE_PLAN, "请帮我系统学习晶体管放大电路")

    assert plan.title == "晶体管放大电路学习规划"
    assert "完整知识链" in plan.goal
    assert len(plan.stages) == 4
    assert plan.stages[0].duration == "1～2 小时"
    assert "直流通路" in plan.stages[0].actions[0]
    assert "$I_{BQ}$" in plan.stages[1].goal
    assert "\\frac{\\beta R'_L}{r_{be}}" in " ".join(plan.stages[2].actions)
    assert "85%" in " ".join(plan.metrics)
    assert len(plan.schedule) == 5

    generic = parse_learning_plan(
        "# 学习规划\n\n目标：掌握反馈放大电路。\n\n## 阶段一：基础\n- 行动：完成概念复习",
        "系统掌握反馈放大电路",
    )
    assert generic.title == "系统掌握反馈放大电路 · 学习规划"


def test_build_learning_plan_ppt_is_editable_and_in_bounds(tmp_path: Path) -> None:
    output = tmp_path / "plan.pptx"
    plan, slide_count = build_learning_plan_ppt(SAMPLE_PLAN, output)

    assert output.stat().st_size > 30_000
    assert slide_count == 10
    presentation = Presentation(str(output))
    assert len(presentation.slides) == slide_count
    all_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert plan.title in all_text
    assert "学习路线一图看懂" in all_text
    assert "静态工作点分析" in all_text
    assert "什么时候算真正学会" in all_text
    assert "I_BQ" not in all_text

    formula_pictures = [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        and shape._element.nvPicPr.cNvPr.get("name") == "LaTeX formula"
    ]
    assert len(formula_pictures) >= 3
    assert any("IBQ" in (shape._element.nvPicPr.cNvPr.get("descr") or "") for shape in formula_pictures)
    with ZipFile(output) as package:
        formula_media = [name for name in package.namelist() if name.startswith("ppt/media/")]
    assert len(formula_media) >= 3

    for slide in presentation.slides:
        for shape in slide.shapes:
            assert shape.left >= 0
            assert shape.top >= 0
            assert shape.left + shape.width <= presentation.slide_width + 2
            assert shape.top + shape.height <= presentation.slide_height + 2


def test_generate_learning_plan_ppt_reuses_content_cache(tmp_path: Path) -> None:
    first_path, first_title, first_count = generate_learning_plan_ppt(
        tmp_path,
        "student-session",
        SAMPLE_PLAN,
        "晶体管放大电路",
    )
    first_mtime = first_path.stat().st_mtime_ns
    second_path, second_title, second_count = generate_learning_plan_ppt(
        tmp_path,
        "student-session",
        SAMPLE_PLAN,
        "晶体管放大电路",
    )

    assert second_path == first_path
    assert second_path.stat().st_mtime_ns == first_mtime
    assert second_title == first_title
    assert second_count == first_count
    assert presentation_filename(first_title).endswith(".pptx")
