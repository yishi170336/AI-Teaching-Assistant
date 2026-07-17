from __future__ import annotations

import hashlib
import logging
import os
import re
import tempfile
from dataclasses import dataclass, field
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


logger = logging.getLogger(__name__)


SLIDE_WIDTH = Inches(13.333333)
SLIDE_HEIGHT = Inches(7.5)
FONT_CN = "Microsoft YaHei"
FONT_LATIN = "Aptos"

INK = "163E3D"
INK_SOFT = "3B5E5B"
TEAL_DARK = "0C5F5B"
TEAL = "0F766E"
TEAL_BRIGHT = "2AA695"
MINT = "DFF2ED"
MINT_LIGHT = "F2F8F6"
CREAM = "FBFAF5"
WHITE = "FFFFFF"
SAND = "EEE8DA"
CORAL = "E98268"
GOLD = "D5A943"
GRAY = "718784"
LINE = "D7E5E1"

_HEADING_RE = re.compile(r"^(#{1,4})\s+(.+?)\s*$", re.MULTILINE)
_LIST_RE = re.compile(r"^\s*(?:[-*+]\s+|\d+[.)、]\s+|[（(]?[一二三四五六七八九十]+[）)、.]\s*)")
_STAGE_RE = re.compile(
    r"(?:阶段\s*[一二三四五六七八九十\d]*|第\s*[一二三四五六七八九十\d]+\s*(?:阶段|步|课|天|周)|"
    r"模块\s*[一二三四五六七八九十\d]+|课次\s*[一二三四五六七八九十\d]+)",
    re.IGNORECASE,
)
_DURATION_RE = re.compile(
    r"\d+(?:\.\d+)?(?:\s*[~～\-–—至]\s*\d+(?:\.\d+)?)?\s*(?:分钟|小时|课时|天|周)",
    re.IGNORECASE,
)
_LATEX_SPAN_RE = re.compile(
    r"\$\$(.+?)\$\$|(?<!\\)\$(.+?)(?<!\\)\$|\\\[(.+?)\\\]|\\\((.+?)\\\)",
    re.DOTALL,
)
_MATH_FONT_REGULAR = Path("C:/Windows/Fonts/msyh.ttc")
_MATH_FONT_BOLD = Path("C:/Windows/Fonts/msyhbd.ttc")


@dataclass
class PlanSection:
    title: str
    body: str
    level: int = 2


@dataclass
class PlanStage:
    title: str
    goal: str
    actions: list[str] = field(default_factory=list)
    standards: list[str] = field(default_factory=list)
    sources: list[str] = field(default_factory=list)
    duration: str = ""


@dataclass
class ParsedLearningPlan:
    title: str
    goal: str
    summary: str
    stages: list[PlanStage]
    schedule: list[str]
    metrics: list[str]
    principles: list[str]


def _mathtext_safe(value: str) -> str:
    value = _normalize_latex(_plain(value))

    def normalize_formula(match: re.Match[str]) -> str:
        body = next((part for part in match.groups() if part is not None), "")
        body = body.replace(r"\displaystyle", "")
        body = re.sub(r"\\begin\{(?:aligned|align\*?|array)\}", "", body)
        body = re.sub(r"\\end\{(?:aligned|align\*?|array)\}", "", body)
        body = body.replace(r"\\", r"\quad ")
        body = re.sub(r"\\operatorname\{([^{}]+)\}", r"\\mathrm{\1}", body)
        body = re.sub(
            r"\\text\{([A-Za-z0-9 ./%°μΩ-]+)\}",
            r"\\mathrm{\1}",
            body,
        )
        return f"${body.strip()}$"

    return _LATEX_SPAN_RE.sub(normalize_formula, value)


def _wrap_rich_text(value: str, capacity: int) -> str:
    value = _mathtext_safe(value)
    capacity = max(6, capacity)
    lines: list[str] = []
    current: list[str] = []
    current_length = 0
    cursor = 0

    def add_plain(chunk: str) -> None:
        nonlocal current, current_length
        for char in chunk:
            if char == "\n":
                lines.append("".join(current).rstrip())
                current = []
                current_length = 0
                continue
            if current_length >= capacity:
                lines.append("".join(current).rstrip())
                current = []
                current_length = 0
            current.append(char)
            current_length += 1

    for match in _LATEX_SPAN_RE.finditer(value):
        add_plain(value[cursor : match.start()])
        formula = match.group(0)
        formula_length = max(3, len(_clean_math(formula)))
        if current and current_length + formula_length > capacity:
            lines.append("".join(current).rstrip())
            current = []
            current_length = 0
        current.append(formula)
        current_length += min(formula_length, capacity)
        cursor = match.end()
    add_plain(value[cursor:])
    if current or not lines:
        lines.append("".join(current).rstrip())
    return "\n".join(line for line in lines if line)


def _fit_rich_text(value: str, width: float, height: float, size: float) -> str:
    capacity = max(7, int(width * 72 / max(size * 0.9, 1)))
    max_lines = max(1, int(height * 72 / max(size * 1.02, 1)))
    fitted = _shorten(value, capacity * max_lines)
    return _wrap_rich_text(fitted, capacity)


class LatexTextRenderer:
    def __init__(self, scratch_dir: Path):
        self.scratch_dir = scratch_dir
        self.scratch_dir.mkdir(parents=True, exist_ok=True)

    @staticmethod
    def _font(bold: bool):
        from matplotlib.font_manager import FontProperties

        font_path = _MATH_FONT_BOLD if bold and _MATH_FONT_BOLD.exists() else _MATH_FONT_REGULAR
        if font_path.exists():
            return FontProperties(fname=str(font_path))
        return FontProperties(family=FONT_CN, weight="bold" if bold else "normal")

    @staticmethod
    def _figure(width: float, height: float):
        from matplotlib.backends.backend_agg import FigureCanvasAgg
        from matplotlib.figure import Figure

        figure = Figure(figsize=(max(width, 0.2), max(height, 0.2)), dpi=200)
        figure.patch.set_alpha(0)
        FigureCanvasAgg(figure)
        return figure

    def render_text(
        self,
        value: str,
        width: float,
        height: float,
        *,
        size: float,
        color: str,
        bold: bool,
        align: Any,
        valign: Any,
        line_spacing: float,
    ) -> Path:
        payload = _fit_rich_text(value, width, height, size)
        key = hashlib.sha256(
            repr(("text", payload, width, height, size, color, bold, str(align), str(valign))).encode("utf-8")
        ).hexdigest()[:24]
        output = self.scratch_dir / f"formula-{key}.png"
        if output.exists():
            return output

        from matplotlib import rc_context

        horizontal = "center" if align == PP_ALIGN.CENTER else "right" if align == PP_ALIGN.RIGHT else "left"
        x = 0.5 if horizontal == "center" else 0.99 if horizontal == "right" else 0.01
        if valign == MSO_ANCHOR.MIDDLE:
            vertical, y = "center", 0.5
        elif valign == MSO_ANCHOR.BOTTOM:
            vertical, y = "bottom", 0.02
        else:
            vertical, y = "top", 0.98
        figure = self._figure(width, height)
        with rc_context(
            {
                "mathtext.fontset": "stix",
                "mathtext.default": "bf" if bold else "it",
                "axes.unicode_minus": False,
            }
        ):
            figure.text(
                x,
                y,
                payload,
                color=f"#{color}",
                fontsize=size,
                fontproperties=self._font(bold),
                horizontalalignment=horizontal,
                verticalalignment=vertical,
                linespacing=line_spacing,
            )
            figure.canvas.print_png(str(output))
        return output

    def render_lines(
        self,
        lines: list[tuple[str, str]],
        width: float,
        height: float,
        *,
        size: float,
        color: str,
        marker_color: str,
    ) -> Path:
        row_height = height / max(len(lines), 1)
        fitted = [
            (marker, _fit_rich_text(value, width - 0.54, row_height, size))
            for marker, value in lines
        ]
        key = hashlib.sha256(
            repr(("lines", fitted, width, height, size, color, marker_color)).encode("utf-8")
        ).hexdigest()[:24]
        output = self.scratch_dir / f"formula-list-{key}.png"
        if output.exists():
            return output

        from matplotlib import rc_context

        figure = self._figure(width, height)
        marker_x = 0.01
        text_x = min(0.14, 0.52 / max(width, 0.1))
        with rc_context({"mathtext.fontset": "stix", "axes.unicode_minus": False}):
            for index, (marker, value) in enumerate(fitted):
                y = 1 - (index + 0.5) / max(len(fitted), 1)
                figure.text(
                    marker_x,
                    y,
                    marker,
                    color=f"#{marker_color}",
                    fontsize=size,
                    fontproperties=self._font(True),
                    verticalalignment="center",
                )
                figure.text(
                    text_x,
                    y,
                    value,
                    color=f"#{color}",
                    fontsize=size,
                    fontproperties=self._font(False),
                    verticalalignment="center",
                    linespacing=1.08,
                )
            figure.canvas.print_png(str(output))
        return output


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _normalize_latex(value: str) -> str:
    def canonical(match: re.Match[str]) -> str:
        body = next((part for part in match.groups() if part is not None), "")
        body = re.sub(r"\s+", " ", body).strip()
        return f"${body}$"

    return _LATEX_SPAN_RE.sub(canonical, value)


def _has_latex(value: str) -> bool:
    return bool(_LATEX_SPAN_RE.search(_normalize_latex(value)))


def _clean_math(value: str) -> str:
    replacements = {
        r"\geq": "≥",
        r"\ge": "≥",
        r"\leq": "≤",
        r"\le": "≤",
        r"\approx": "≈",
        r"\times": "×",
        r"\cdot": "·",
        r"\beta": "β",
        r"\alpha": "α",
        r"\Delta": "Δ",
        r"\Omega": "Ω",
        r"\mu": "μ",
        r"\text": "",
    }
    for source, target in replacements.items():
        value = value.replace(source, target)
    value = re.sub(r"\\frac\s*\{([^{}]+)\}\s*\{([^{}]+)\}", r"\1/\2", value)
    value = re.sub(r"_\{([^{}]+)\}", r"\1", value)
    value = re.sub(r"\^\{([^{}]+)\}", r"\1", value)
    value = value.replace("{", "").replace("}", "").replace("$", "")
    return value


def _plain(value: str, preserve_latex: bool = True) -> str:
    value = _normalize_latex(value)
    formulas: list[str] = []

    def protect(match: re.Match[str]) -> str:
        formulas.append(match.group(0))
        return f"ZXQLATEX{len(formulas) - 1}QXZ"

    value = _LATEX_SPAN_RE.sub(protect, value)
    value = re.sub(r"!\[([^]]*)\]\([^)]+\)", r"\1", value)
    value = re.sub(r"\[([^]]+)\]\([^)]+\)", r"\1", value)
    value = re.sub(r"<[^>]+>", "", value)
    value = value.replace("**", "").replace("__", "").replace("`", "")
    value = _LIST_RE.sub("", value.strip())
    value = re.sub(r"\s+", " ", value).strip(" \t\r\n|—–-：:")
    for index, formula in enumerate(formulas):
        replacement = formula if preserve_latex else _clean_math(formula)
        value = value.replace(f"ZXQLATEX{index}QXZ", replacement)
    return value


def _shorten(value: str, limit: int) -> str:
    value = _plain(value)
    visible = _plain(value, preserve_latex=False)
    if len(visible) <= limit:
        return value
    content_limit = max(1, limit - 1)
    result: list[str] = []
    consumed = 0
    cursor = 0
    for match in _LATEX_SPAN_RE.finditer(value):
        prefix = value[cursor : match.start()]
        remaining = content_limit - consumed
        if remaining <= 0:
            break
        if len(prefix) > remaining:
            result.append(prefix[:remaining])
            consumed = content_limit
            break
        result.append(prefix)
        consumed += len(prefix)
        formula = match.group(0)
        formula_length = max(3, len(_clean_math(formula)))
        if consumed + formula_length > content_limit:
            break
        result.append(formula)
        consumed += formula_length
        cursor = match.end()
    else:
        remaining = content_limit - consumed
        if remaining > 0:
            result.append(value[cursor : cursor + remaining])
    clipped = "".join(result).rstrip("，、；：,. ")
    return f"{clipped}…"


def _dedupe(values: list[str], limit: int = 99) -> list[str]:
    result: list[str] = []
    seen: set[str] = set()
    for value in values:
        cleaned = _plain(value)
        key = re.sub(r"[\s，。；、,.]", "", _plain(cleaned, preserve_latex=False)).lower()
        if not cleaned or key in seen:
            continue
        seen.add(key)
        result.append(cleaned)
        if len(result) >= limit:
            break
    return result


def _split_sections(markdown: str) -> tuple[str, list[PlanSection]]:
    matches = list(_HEADING_RE.finditer(markdown))
    if not matches:
        return "", [PlanSection(title="学习方案", body=markdown, level=2)]
    preamble = markdown[: matches[0].start()].strip()
    sections: list[PlanSection] = []
    for index, match in enumerate(matches):
        end = matches[index + 1].start() if index + 1 < len(matches) else len(markdown)
        sections.append(
            PlanSection(
                title=_plain(match.group(2)),
                body=markdown[match.end() : end].strip(),
                level=len(match.group(1)),
            )
        )
    return preamble, sections


def _content_items(value: str, limit: int = 12) -> list[str]:
    items: list[str] = []
    for raw_line in value.splitlines():
        line = raw_line.strip()
        if not line or line.startswith("#") or re.fullmatch(r"[-*_]{3,}", line):
            continue
        if line.startswith("|") and line.endswith("|"):
            cells = [_plain(cell) for cell in line.strip("|").split("|")]
            if cells and not all(re.fullmatch(r":?-{2,}:?", cell or "--") for cell in cells):
                line = " · ".join(cell for cell in cells if cell)
            else:
                continue
        cleaned = _plain(line)
        if cleaned:
            items.append(cleaned)
    if len(items) <= 1:
        prose = _plain(value)
        items.extend(
            part.strip()
            for part in re.split(r"(?<=[。；！？])", prose)
            if part.strip()
        )
    return _dedupe(items, limit)


def _strip_label(value: str) -> tuple[str, str]:
    match = re.match(
        r"^(目标|阶段目标|学习目标|行动|具体行动|学习任务|任务|练习|完成标准|验收标准|"
        r"完成标志|资料依据|依据|参考资料|建议时长|预计用时|时间)\s*[:：]?\s*(.*)$",
        value,
    )
    if not match:
        return "", value
    return match.group(1), match.group(2).strip()


def _stage_details(section: PlanSection) -> PlanStage:
    buckets: dict[str, list[str]] = {
        "goal": [],
        "actions": [],
        "standards": [],
        "sources": [],
    }
    current = ""
    general: list[str] = []
    for item in _content_items(section.body, 28):
        label, payload = _strip_label(item)
        if label:
            if "目标" in label:
                current = "goal"
            elif label in {"行动", "具体行动", "学习任务", "任务", "练习"}:
                current = "actions"
            elif "标准" in label or "标志" in label:
                current = "standards"
            elif "依据" in label or "资料" in label:
                current = "sources"
            elif "时长" in label or "用时" in label or label == "时间":
                current = ""
            if payload and current:
                buckets[current].append(payload)
            continue
        if re.search(r"\[资料\s*\d+\]", item) or item.startswith("资料"):
            buckets["sources"].append(item)
        elif current:
            buckets[current].append(item)
        else:
            general.append(item)

    title = re.sub(r"^[一二三四五六七八九十\d]+[、.．]\s*", "", section.title)
    title = re.sub(
        r"^(?:阶段\s*[一二三四五六七八九十\d]+|第\s*[一二三四五六七八九十\d]+\s*阶段)\s*[:：·|-]?\s*",
        "",
        title,
    )
    title = _DURATION_RE.sub("", title)
    title = re.sub(r"[（(]\s*[）)]", "", title).strip(" ：:·|-—（）()")
    goal_candidates = buckets["goal"] or general[:1]
    actions = buckets["actions"] or general[1:] or general
    standards = buckets["standards"]
    if not standards:
        standards = [item for item in general if re.search(r"完成|正确率|能够|独立|通过|达到", item)]
    duration_match = _DURATION_RE.search(f"{section.title} {section.body[:500]}")
    return PlanStage(
        title=_shorten(title, 32) or "学习任务",
        goal=_shorten(goal_candidates[0] if goal_candidates else title, 72),
        actions=[_shorten(item, 78) for item in _dedupe(actions, 5)],
        standards=[_shorten(item, 64) for item in _dedupe(standards, 3)],
        sources=[_shorten(item, 96) for item in _dedupe(buckets["sources"], 2)],
        duration=duration_match.group(0) if duration_match else "",
    )


def _generic_title(value: str) -> bool:
    compact = re.sub(r"[\s：:·|—-]", "", value)
    return not compact or compact in {
        "学习规划",
        "学习计划",
        "学习规划方案",
        "个性化学习规划",
        "学习路线",
        "学习方案",
    }


def _derive_title(sections: list[PlanSection], goal: str, topic: str) -> str:
    for section in sections[:3]:
        if _STAGE_RE.search(section.title) or re.search(r"日程|安排|验收|指标", section.title):
            continue
        if section.level <= 2 and not _generic_title(section.title):
            title = re.sub(r"^(学习规划|学习计划|学习路线)\s*[:：·|-]\s*", "", section.title)
            if title:
                return _shorten(title, 26)
    topic_text = _plain(topic)
    if topic_text:
        topic_text = re.sub(r"^(请|帮我|请帮我|我想|我需要)", "", topic_text)
        topic_text = re.sub(r"(制定|生成|做)(一份|一个)?(详细的)?(学习规划|学习计划).*$", "", topic_text)
        if topic_text:
            return _shorten(topic_text, 26)
    return _shorten(goal, 26) or "个性化学习规划"


def parse_learning_plan(markdown: str, topic: str = "") -> ParsedLearningPlan:
    preamble, sections = _split_sections(markdown)
    relevant = [
        section
        for section in sections
        if not re.search(r"检索依据|参考文献|引用资料|sources?", section.title, re.IGNORECASE)
    ]
    stage_sections = [section for section in relevant if _STAGE_RE.search(section.title)]
    if not stage_sections:
        stage_sections = [
            section
            for section in relevant
            if re.search(r"目标|行动|任务|完成标准|学习内容|练习", section.body)
            and not re.search(r"总体|验收|指标|日程|安排|时间表", section.title)
        ]
    if not stage_sections:
        stage_sections = [section for section in relevant if section.body][:5]

    stages = [_stage_details(section) for section in stage_sections[:6]]
    stages = [stage for stage in stages if stage.title or stage.actions]
    if not stages:
        fallback = _content_items(markdown, 8)
        stages = [
            PlanStage(
                title="核心学习任务",
                goal=_shorten(fallback[0] if fallback else "完成本次学习目标", 95),
                actions=[_shorten(item, 88) for item in fallback[1:6]],
                standards=["能独立复述核心概念，并完成一次自测与纠错"],
            )
        ]

    goal_candidates: list[str] = []
    searchable = "\n".join([preamble] + [section.body for section in relevant[:3]])
    goal_match = re.search(
        r"(?:总体目标|学习目标|规划目标|目标)\s*[:：]\s*([^\n]{4,180})",
        searchable,
    )
    if goal_match:
        goal_candidates.append(goal_match.group(1))
    goal_candidates.extend(_content_items(preamble, 2))
    goal_candidates.extend(stage.goal for stage in stages[:1])
    goal = _shorten(next((item for item in goal_candidates if _plain(item)), "完成本次学习目标"), 82)

    schedule_sections = [
        section
        for section in relevant
        if re.search(r"日程|安排|时间表|每日|每周|进度|Day\s*\d+", section.title, re.IGNORECASE)
    ]
    schedule = _dedupe(
        [item for section in schedule_sections for item in _content_items(section.body, 12)],
        7,
    )

    metric_sections = [
        section
        for section in relevant
        if re.search(r"验收|指标|完成标准|自测|复盘|检查", section.title)
        and section not in stage_sections
    ]
    metrics = [item for section in metric_sections for item in _content_items(section.body, 12)]
    metrics.extend(item for stage in stages for item in stage.standards)
    metrics = [_shorten(item, 82) for item in _dedupe(metrics, 6)]
    if not metrics:
        metrics = [
            "能不看资料复述核心概念与关键关系",
            "完成代表性练习并解释每一步依据",
            "记录错因，完成一次针对性复练",
            "用自测结果决定是否进入下一阶段",
        ]

    overview_sections = [
        section
        for section in relevant
        if re.search(r"评估|诊断|原则|说明|节奏|现状", section.title)
        and section not in stage_sections
    ]
    principles = [item for section in overview_sections for item in _content_items(section.body, 8)]
    principles = [_shorten(item, 76) for item in _dedupe(principles, 4)]
    default_principles = [
        "难点用例题与错题双向验证",
        "根据自测结果动态调整学习强度",
        "练习、纠错、复盘形成闭环",
    ]
    principles = _dedupe(principles + default_principles, 3)

    summary_candidates = _content_items(preamble, 2)
    if not summary_candidates:
        for section in relevant:
            if section not in stage_sections:
                summary_candidates.extend(_content_items(section.body, 1))
                if summary_candidates:
                    break
    summary = _shorten(summary_candidates[0] if summary_candidates else goal, 150)
    title = _derive_title(sections, goal, topic)
    if not title.endswith(("规划", "计划", "路线")):
        title = f"{title} · 学习规划"

    return ParsedLearningPlan(
        title=title,
        goal=goal,
        summary=summary,
        stages=stages,
        schedule=[_shorten(item, 94) for item in schedule],
        metrics=metrics,
        principles=principles,
    )


def _set_background(slide, color: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(color)


def _shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str,
    line: str | None = None,
    line_width: float = 1,
    radius_adjust: float | None = None,
):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line or fill)
    shape.line.width = Pt(line_width)
    if radius_adjust is not None and shape.adjustments:
        shape.adjustments[0] = radius_adjust
    return shape


def _line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = LINE,
    width: float = 1.4,
    dash: bool = False,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = _rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def _text(
    slide,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    font: str = FONT_CN,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0,
    line_spacing: float = 1.08,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _picture_alt_text(picture, value: str) -> None:
    try:
        properties = picture._element.nvPicPr.cNvPr
        properties.set("name", "LaTeX formula")
        properties.set("descr", f"LaTeX rendered formula: {_plain(value, preserve_latex=False)}")
    except (AttributeError, TypeError):
        logger.debug("Unable to attach alternative text to rendered formula", exc_info=True)


def _display_text(
    slide,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    renderer: LatexTextRenderer | None,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    font: str = FONT_CN,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0,
    line_spacing: float = 1.08,
):
    if renderer is not None and _has_latex(value):
        try:
            image_path = renderer.render_text(
                value,
                w,
                h,
                size=size,
                color=color,
                bold=bold,
                align=align,
                valign=valign,
                line_spacing=line_spacing,
            )
            picture = slide.shapes.add_picture(
                str(image_path),
                Inches(x),
                Inches(y),
                width=Inches(w),
                height=Inches(h),
            )
            _picture_alt_text(picture, value)
            return picture
        except Exception:
            logger.warning("LaTeX rendering failed; using readable text fallback", exc_info=True)
            value = _plain(value, preserve_latex=False)
    return _text(
        slide,
        value,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        font=font,
        align=align,
        valign=valign,
        margin=margin,
        line_spacing=line_spacing,
    )


def _rich_lines(
    slide,
    lines: list[tuple[str, str]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    renderer: LatexTextRenderer | None = None,
    size: float = 17,
    color: str = INK_SOFT,
    bullet_color: str = TEAL,
    gap: float = 8,
):
    if renderer is not None and any(_has_latex(value) for _, value in lines):
        try:
            image_path = renderer.render_lines(
                lines,
                w,
                h,
                size=size,
                color=color,
                marker_color=bullet_color,
            )
            picture = slide.shapes.add_picture(
                str(image_path),
                Inches(x),
                Inches(y),
                width=Inches(w),
                height=Inches(h),
            )
            _picture_alt_text(picture, " / ".join(value for _, value in lines))
            return picture
        except Exception:
            logger.warning("LaTeX list rendering failed; using readable text fallback", exc_info=True)
            lines = [(marker, _plain(value, preserve_latex=False)) for marker, value in lines]
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0)
    frame.margin_top = frame.margin_bottom = Inches(0)
    for index, (marker, value) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.12
        marker_run = paragraph.add_run()
        marker_run.text = f"{marker}  "
        marker_run.font.name = FONT_LATIN
        marker_run.font.size = Pt(size)
        marker_run.font.bold = True
        marker_run.font.color.rgb = _rgb(bullet_color)
        text_run = paragraph.add_run()
        text_run.text = value
        text_run.font.name = FONT_CN
        text_run.font.size = Pt(size)
        text_run.font.color.rgb = _rgb(color)
    return box


def _brand(slide, *, dark: bool = False) -> None:
    color = "A9CFC8" if dark else TEAL
    _shape(slide, MSO_SHAPE.OVAL, 0.72, 0.46, 0.16, 0.16, fill=CORAL, line=CORAL)
    _text(
        slide,
        "CIRCUITMIND  /  LEARNING LAB",
        0.98,
        0.42,
        3.5,
        0.25,
        size=9,
        color=color,
        bold=True,
        font=FONT_LATIN,
    )


def _footer(slide, page_number: int, *, dark: bool = False) -> None:
    color = "8DB5AF" if dark else "8CA09D"
    _text(slide, "由本次学习规划自动生成 · 内容与图形均可编辑", 0.72, 7.08, 5.2, 0.2, size=8, color=color)
    _text(
        slide,
        f"{page_number:02d}",
        12.0,
        7.02,
        0.58,
        0.24,
        size=9,
        color=color,
        bold=True,
        font=FONT_LATIN,
        align=PP_ALIGN.RIGHT,
    )


def _slide_title(slide, kicker: str, title: str, page_number: int, *, dark: bool = False) -> None:
    _brand(slide, dark=dark)
    _text(slide, kicker.upper(), 0.72, 0.98, 3.2, 0.24, size=10, color=CORAL, bold=True, font=FONT_LATIN)
    _text(slide, title, 0.72, 1.25, 11.7, 0.58, size=30, color=WHITE if dark else INK, bold=True)
    _footer(slide, page_number, dark=dark)


def _add_circuit_art(slide) -> None:
    _shape(slide, MSO_SHAPE.OVAL, 8.85, 1.05, 3.65, 3.65, fill="154E4B", line="2B7770", line_width=1.2)
    _shape(slide, MSO_SHAPE.OVAL, 9.35, 1.55, 2.65, 2.65, fill=TEAL_DARK, line="3D8E84", line_width=1.0)
    _shape(slide, MSO_SHAPE.OVAL, 10.12, 2.32, 1.12, 1.12, fill=CORAL, line=CORAL)
    _text(slide, "PLAN", 10.18, 2.69, 1.0, 0.22, size=10, color=WHITE, bold=True, font=FONT_LATIN, align=PP_ALIGN.CENTER)
    nodes = [(8.2, 1.7), (8.1, 4.4), (10.0, 5.35), (12.1, 4.8), (12.55, 2.05)]
    center = (10.68, 2.88)
    for x, y in nodes:
        _line(slide, center[0], center[1], x + 0.1, y + 0.1, color="3D8E84", width=1.2)
        _shape(slide, MSO_SHAPE.OVAL, x, y, 0.2, 0.2, fill=GOLD if y > 4 else MINT, line=GOLD if y > 4 else MINT)
    for row in range(3):
        for col in range(4):
            _shape(slide, MSO_SHAPE.OVAL, 8.05 + col * 0.42, 5.85 + row * 0.34, 0.055, 0.055, fill="4F8C85", line="4F8C85")


def _title_slide(prs: Presentation, plan: ParsedLearningPlan, renderer: LatexTextRenderer) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, TEAL_DARK)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.18, 7.5, fill=CORAL, line=CORAL)
    _brand(slide, dark=True)
    _text(slide, "PERSONAL LEARNING BLUEPRINT", 0.78, 1.35, 5.0, 0.28, size=11, color="8CCDC3", bold=True, font=FONT_LATIN)
    _text(slide, plan.title, 0.75, 1.83, 7.3, 1.5, size=38, color=WHITE, bold=True)
    _display_text(slide, plan.goal, 0.78, 3.52, 6.55, 1.15, renderer=renderer, size=18, color="D4EAE6")
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.78, 5.2, 2.02, 0.5, fill="164E4A", line="377C75")
    _text(slide, f"{len(plan.stages)} 个执行阶段", 0.78, 5.2, 2.02, 0.5, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.98, 5.2, 2.26, 0.5, fill="164E4A", line="377C75")
    _text(slide, "含量化验收闭环", 2.98, 5.2, 2.26, 0.5, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _text(slide, date.today().strftime("%Y.%m.%d"), 0.8, 6.76, 1.5, 0.24, size=9, color="86AAA5", bold=True, font=FONT_LATIN)
    _add_circuit_art(slide)


def _overview_slide(
    prs: Presentation,
    plan: ParsedLearningPlan,
    page: int,
    renderer: LatexTextRenderer,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, CREAM)
    _slide_title(slide, "01 / NORTH STAR", "先看全局：这份规划要解决什么", page)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 2.02, 7.2, 2.25, fill=TEAL_DARK, line=TEAL_DARK)
    _text(slide, "核心目标", 1.02, 2.35, 1.4, 0.28, size=12, color="8FCFC6", bold=True)
    _display_text(slide, plan.goal, 1.02, 2.86, 6.2, 1.05, renderer=renderer, size=23, color=WHITE, bold=True)

    facts = [
        (f"{len(plan.stages):02d}", "执行阶段"),
        (f"{len(plan.metrics):02d}", "验收指标"),
        ("闭环", "学习方式"),
    ]
    for index, (value, label) in enumerate(facts):
        x = 8.25 + (index % 2) * 2.12
        y = 2.02 + (index // 2) * 1.2
        w = 1.86 if index < 2 else 3.98
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 1.02, fill=WHITE, line=LINE)
        _text(slide, value, x + 0.2, y + 0.18, w - 0.4, 0.36, size=22, color=TEAL, bold=True, font=FONT_LATIN)
        _text(slide, label, x + 0.2, y + 0.64, w - 0.4, 0.2, size=10, color=GRAY, bold=True)

    _text(slide, "执行原则", 0.78, 4.76, 1.5, 0.3, size=13, color=TEAL, bold=True)
    principles = plan.principles[:3]
    card_width = 3.87
    for index, principle in enumerate(principles):
        x = 0.72 + index * 4.13
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 5.2, card_width, 1.25, fill=WHITE, line=LINE)
        _text(slide, f"0{index + 1}", x + 0.22, 5.43, 0.48, 0.3, size=15, color=CORAL, bold=True, font=FONT_LATIN)
        _display_text(slide, principle, x + 0.78, 5.36, 2.82, 0.7, renderer=renderer, size=14, color=INK_SOFT, bold=True)


def _roadmap_slide(
    prs: Presentation,
    plan: ParsedLearningPlan,
    page: int,
    renderer: LatexTextRenderer,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, MINT_LIGHT)
    _slide_title(slide, "02 / LEARNING PATH", "从起点到验收：学习路线一图看懂", page)
    stages = plan.stages
    rows = 1 if len(stages) <= 4 else 2
    per_row = len(stages) if rows == 1 else (len(stages) + 1) // 2
    for row in range(rows):
        row_stages = stages[row * per_row : (row + 1) * per_row]
        if not row_stages:
            continue
        gap = 0.28
        available = 11.86
        width = min(2.76, (available - gap * (len(row_stages) - 1)) / len(row_stages))
        total = width * len(row_stages) + gap * (len(row_stages) - 1)
        start_x = (13.333 - total) / 2
        y = 2.2 + row * 2.15
        for index, stage in enumerate(row_stages):
            global_index = row * per_row + index
            x = start_x + index * (width + gap)
            if index < len(row_stages) - 1:
                _line(slide, x + width, y + 0.72, x + width + gap, y + 0.72, color="8DC9C0", width=2.2)
            _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, 1.66, fill=WHITE, line="CFE3DE", line_width=1.1)
            _shape(slide, MSO_SHAPE.OVAL, x + 0.18, y + 0.18, 0.5, 0.5, fill=CORAL if global_index == 0 else TEAL, line=CORAL if global_index == 0 else TEAL)
            _text(slide, str(global_index + 1), x + 0.18, y + 0.18, 0.5, 0.5, size=13, color=WHITE, bold=True, font=FONT_LATIN, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            _text(slide, _shorten(stage.title, 22), x + 0.82, y + 0.18, width - 1.0, 0.58, size=15, color=INK, bold=True)
            _display_text(
                slide,
                _shorten(stage.goal, 44),
                x + 0.2,
                y + 0.92,
                width - 0.4,
                0.52,
                renderer=renderer,
                size=11,
                color=INK_SOFT,
            )
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.33, 6.23, 6.67, 0.46, fill=MINT, line=MINT)
    _text(slide, "推进规则：达到本阶段完成标准后，再进入下一阶段", 3.33, 6.23, 6.67, 0.46, size=11, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def _stage_slide(
    prs: Presentation,
    stage: PlanStage,
    index: int,
    total: int,
    page: int,
    renderer: LatexTextRenderer,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, CREAM if index % 2 else MINT_LIGHT)
    _brand(slide)
    _footer(slide, page)
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 1.03, 0.08, 1.36, fill=CORAL if index == 1 else TEAL, line=CORAL if index == 1 else TEAL)
    _text(slide, f"STAGE {index:02d}  /  {total:02d}", 1.02, 1.02, 2.2, 0.26, size=10, color=CORAL, bold=True, font=FONT_LATIN)
    _text(slide, stage.title, 1.0, 1.34, 8.9, 0.56, size=29, color=INK, bold=True)
    chip = stage.duration or "按完成标准推进"
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.24, 1.25, 2.28, 0.48, fill=MINT, line=MINT)
    _text(slide, chip, 10.24, 1.25, 2.28, 0.48, size=10, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, 2.08, 11.52, 1.12, fill=TEAL_DARK, line=TEAL_DARK)
    _text(slide, "阶段目标", 1.3, 2.32, 1.05, 0.25, size=11, color="8ECBC2", bold=True)
    _display_text(
        slide,
        stage.goal,
        2.42,
        2.22,
        9.55,
        0.66,
        renderer=renderer,
        size=19,
        color=WHITE,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )

    _text(slide, "要做什么", 1.02, 3.58, 2.0, 0.3, size=14, color=TEAL, bold=True)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0, 4.0, 7.45, 2.38, fill=WHITE, line=LINE)
    actions = stage.actions or ["阅读并整理本阶段核心概念", "完成代表性练习并记录错因", "对照完成标准进行自测"]
    _rich_lines(
        slide,
        [
            (f"{idx + 1:02d}", _shorten(action, 52 if len(actions) >= 5 else 66))
            for idx, action in enumerate(actions[:5])
        ],
        1.3,
        4.28,
        6.85,
        1.86,
        renderer=renderer,
        size=14.5 if len(actions) >= 5 else 16,
        gap=5,
    )

    _text(slide, "如何算完成", 8.78, 3.58, 2.0, 0.3, size=14, color=TEAL, bold=True)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.75, 4.0, 3.77, 2.38, fill="E6F3EF", line="C9E2DC")
    standards = stage.standards or ["能独立讲清核心概念", "完成一次自测并订正错误"]
    _rich_lines(
        slide,
        [("✓", _shorten(standard, 50)) for standard in standards[:3]],
        9.03,
        4.31,
        3.2,
        1.7,
        renderer=renderer,
        size=14,
        color=INK_SOFT,
        bullet_color=TEAL,
        gap=9,
    )
    if stage.sources:
        source_line = _shorten(" / ".join(stage.sources), 82)
        _text(slide, f"资料依据  ·  {source_line}", 1.02, 6.63, 10.9, 0.22, size=8.5, color=GRAY)


def _schedule_slide(
    prs: Presentation,
    plan: ParsedLearningPlan,
    page: int,
    renderer: LatexTextRenderer,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, CREAM)
    _slide_title(slide, "SCHEDULE / CADENCE", "把路线放进时间：按节奏持续推进", page)
    items = plan.schedule[:6]
    _line(slide, 1.35, 2.15, 1.35, 6.43, color="B6D8D2", width=2.2)
    for index, item in enumerate(items):
        y = 2.08 + index * (4.25 / max(len(items), 1))
        _shape(slide, MSO_SHAPE.OVAL, 1.12, y, 0.46, 0.46, fill=CORAL if index == 0 else TEAL, line=CREAM, line_width=2.5)
        _text(slide, f"{index + 1:02d}", 1.78, y - 0.03, 0.56, 0.26, size=11, color=CORAL, bold=True, font=FONT_LATIN)
        _display_text(
            slide,
            _shorten(item, 78),
            2.42,
            y - 0.09,
            9.5,
            0.55,
            renderer=renderer,
            size=16,
            color=INK_SOFT,
            bold=True,
        )
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 9.76, 6.18, 2.5, 0.5, fill=TEAL_DARK, line=TEAL_DARK)
    _text(slide, "完成一项，复盘一次", 9.76, 6.18, 2.5, 0.5, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def _metrics_slide(
    prs: Presentation,
    plan: ParsedLearningPlan,
    page: int,
    renderer: LatexTextRenderer,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, TEAL_DARK)
    _slide_title(slide, "QUALITY GATE", "用结果验收：什么时候算真正学会", page, dark=True)
    metrics = plan.metrics[:6]
    columns = 3 if len(metrics) > 4 else 2
    rows = (len(metrics) + columns - 1) // columns
    card_w = 3.72 if columns == 3 else 5.72
    gap_x = 0.3
    start_x = (13.333 - (card_w * columns + gap_x * (columns - 1))) / 2
    card_h = 1.55 if rows <= 2 else 1.22
    gap_y = 0.28
    start_y = 2.2
    for index, metric in enumerate(metrics):
        col = index % columns
        row = index // columns
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h, fill="124B48", line="2D716B")
        _text(slide, f"0{index + 1}", x + 0.25, y + 0.22, 0.58, 0.3, size=15, color=CORAL, bold=True, font=FONT_LATIN)
        metric_limit = 58 if columns == 3 else 76
        _display_text(
            slide,
            _shorten(metric, metric_limit),
            x + 0.96,
            y + 0.2,
            card_w - 1.24,
            card_h - 0.35,
            renderer=renderer,
            size=14.5,
            color=WHITE,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.76, 6.24, 7.82, 0.48, fill="D6EEE9", line="D6EEE9")
    _text(slide, "全部达标 → 进入综合复盘；任一未达标 → 回到对应阶段补强", 2.76, 6.24, 7.82, 0.48, size=11, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def _closing_slide(
    prs: Presentation,
    plan: ParsedLearningPlan,
    page: int,
    renderer: LatexTextRenderer,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, CREAM)
    _brand(slide)
    _footer(slide, page)
    _text(slide, "START SMALL. CLOSE THE LOOP.", 0.78, 1.05, 5.0, 0.3, size=11, color=CORAL, bold=True, font=FONT_LATIN)
    _text(slide, "现在就开始", 0.75, 1.56, 5.8, 0.8, size=38, color=INK, bold=True)
    _text(slide, "先完成最小学习闭环，再逐步加量。", 0.78, 2.48, 5.9, 0.42, size=17, color=GRAY)
    first_actions = plan.stages[0].actions[:3] if plan.stages else []
    if not first_actions:
        first_actions = ["明确本阶段目标", "完成一次核心学习任务", "对照标准完成自测"]
    for index, action in enumerate(first_actions):
        y = 3.3 + index * 0.92
        _shape(slide, MSO_SHAPE.OVAL, 0.8, y, 0.54, 0.54, fill=CORAL if index == 0 else TEAL, line=CORAL if index == 0 else TEAL)
        _text(slide, str(index + 1), 0.8, y, 0.54, 0.54, size=13, color=WHITE, bold=True, font=FONT_LATIN, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _display_text(
            slide,
            _shorten(action, 54),
            1.62,
            y - 0.02,
            5.48,
            0.58,
            renderer=renderer,
            size=16,
            color=INK_SOFT,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )

    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.1, 1.23, 4.23, 4.92, fill=TEAL_DARK, line=TEAL_DARK)
    _text(slide, "今日启动卡", 8.52, 1.72, 2.6, 0.38, size=14, color="8FCFC6", bold=True)
    _text(slide, "01", 8.48, 2.28, 2.2, 1.0, size=58, color=CORAL, bold=True, font=FONT_LATIN)
    _text(slide, plan.stages[0].title if plan.stages else "核心学习任务", 8.54, 3.42, 3.2, 0.8, size=24, color=WHITE, bold=True)
    _line(slide, 8.52, 4.55, 11.82, 4.55, color="3A7772", width=1.0)
    _text(slide, "完成后：记录结果 · 标记疑点 · 决定下一步", 8.52, 4.9, 3.22, 0.72, size=13, color="D2E7E3")
    _text(slide, "把规划变成行动，才是学习真正发生的时刻。", 7.95, 6.54, 4.55, 0.34, size=11, color=TEAL, bold=True, align=PP_ALIGN.RIGHT)


def build_learning_plan_ppt(markdown: str, output_path: Path, topic: str = "") -> tuple[ParsedLearningPlan, int]:
    plan = parse_learning_plan(markdown, topic)
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    core = presentation.core_properties
    core.title = plan.title
    core.subject = "学习规划演示文稿"
    core.author = "CircuitMind"
    core.comments = "根据本次学习规划回答自动生成"

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="circuitmind-latex-") as latex_dir:
        renderer = LatexTextRenderer(Path(latex_dir))
        _title_slide(presentation, plan, renderer)
        page = 2
        _overview_slide(presentation, plan, page, renderer)
        page += 1
        _roadmap_slide(presentation, plan, page, renderer)
        page += 1
        for index, stage in enumerate(plan.stages, start=1):
            _stage_slide(presentation, stage, index, len(plan.stages), page, renderer)
            page += 1
        if plan.schedule:
            _schedule_slide(presentation, plan, page, renderer)
            page += 1
        _metrics_slide(presentation, plan, page, renderer)
        page += 1
        _closing_slide(presentation, plan, page, renderer)
        presentation.save(str(output_path))

    return plan, len(presentation.slides)


def learning_plan_ppt_path(root_dir: Path, session_id: str, markdown: str, topic: str = "") -> Path:
    digest = hashlib.sha256(f"v2-latex\0{topic}\0{markdown}".encode("utf-8")).hexdigest()[:20]
    return root_dir / "data" / "presentations" / session_id / f"learning-plan-{digest}.pptx"


def generate_learning_plan_ppt(
    root_dir: Path,
    session_id: str,
    markdown: str,
    topic: str = "",
) -> tuple[Path, str, int]:
    output_path = learning_plan_ppt_path(root_dir, session_id, markdown, topic)
    plan = parse_learning_plan(markdown, topic)
    if output_path.exists() and output_path.stat().st_size > 10_000:
        slide_count = len(Presentation(str(output_path)).slides)
        return output_path, plan.title, slide_count

    output_path.parent.mkdir(parents=True, exist_ok=True)
    handle, temporary_name = tempfile.mkstemp(
        prefix="learning-plan-",
        suffix=".pptx",
        dir=output_path.parent,
    )
    os.close(handle)
    temporary_path = Path(temporary_name)
    try:
        _, slide_count = build_learning_plan_ppt(markdown, temporary_path, topic)
        os.replace(temporary_path, output_path)
    finally:
        temporary_path.unlink(missing_ok=True)
    return output_path, plan.title, slide_count


def presentation_filename(title: str) -> str:
    safe = re.sub(
        r"[\\/:*?\"<>|\r\n]+",
        "-",
        _plain(title, preserve_latex=False),
    ).strip(" .-")
    safe = _shorten(safe, 48) or "学习规划"
    return f"{safe}.pptx"
